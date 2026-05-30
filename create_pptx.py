"""
PPT 생성 스크립트
서울시 전기차 충전 인프라 분석 발표용
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as In
import sys

OUT_FILE = "서울시_EV_인프라분석_발표_v2.pptx"

# ─── Color palette ────────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C = {
    "darkBg":    "060E1F",
    "midBg":     "0D1B35",
    "cardBg":    "1A2B4A",
    "contentBg": "EEF2FF",
    "accent1":   "00C8FF",
    "accent2":   "3B82F6",
    "accent3":   "10B981",
    "warn":      "F59E0B",
    "danger":    "EF4444",
    "white":     "FFFFFF",
    "offWhite":  "D0DCF0",
    "muted":     "7A94BF",
    "dark":      "1E2B42",
    "tagBg":     "142038",
    "pink":      "F472B6",
}

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

W = 13.33   # slide width in inches
H = 7.5     # slide height in inches

BLANK_LAYOUT = prs.slide_layouts[6]  # blank layout


def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width=0):
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex and line_width > 0:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, italic=False,
             color="FFFFFF", align="left", valign="middle",
             font_name="맑은 고딕"):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = 3   # MSO_ANCHOR.MIDDLE
    elif valign == "top":
        tf.vertical_anchor = 1
    else:
        tf.vertical_anchor = 4   # BOTTOM

    # Remove default margin
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    run.font.name = font_name
    return txBox


def add_title_bar(slide, title, subtitle=""):
    """Dark navy header bar with cyan accent line."""
    rect(slide, 0, 0, W, 1.3, C["midBg"])
    rect(slide, 0, 1.3, W, 0.07, C["accent1"])
    add_text(slide, title, 0.5, 0.12, W - 1, 0.85,
             font_size=30, bold=True, color=C["white"])
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.9, W - 1, 0.38,
                 font_size=12, color=C["offWhite"])


def page_num(slide, n, total):
    add_text(slide, f"{n} / {total}", W - 1.2, H - 0.45, 0.9, 0.3,
             font_size=9, color=C["muted"], align="right")


def section_pill(slide, label, color_hex):
    rect(slide, 0.5, 0.3, 2.0, 0.42, C["cardBg"], color_hex, 1.0)
    add_text(slide, label, 0.5, 0.3, 2.0, 0.42,
             font_size=9, bold=True, color=color_hex, align="center")


def bullet_item(slide, text, x, y, w, h, bullet_color, text_color, font_size=11.5):
    add_text(slide, f"▪  {text}", x, y, w, h,
             font_size=font_size, color=text_color)


TOTAL = 13

# ════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["darkBg"])

rect(s, 0, 0, 0.35, H, C["accent1"])

# Decorative glow circles
from pptx.util import Inches as I_
from pptx.dml.color import RGBColor

circle = s.shapes.add_shape(9, I_(8.5), I_(-0.8), I_(5.5), I_(5.5))  # 9=OVAL
circle.fill.solid()
circle.fill.fore_color.rgb = rgb("001A3A")
circle.line.fill.background()

add_text(s, "⚡", 0.55, 0.6, 0.9, 0.9, font_size=38, color=C["accent1"])

add_text(s, "서울시 전기차 충전 인프라", 0.6, 1.55, W - 1.0, 1.1,
         font_size=42, bold=True, color=C["white"])
add_text(s, "부족 지역 분석 및 시각화 서비스", 0.6, 2.6, W - 1.0, 0.9,
         font_size=34, bold=True, color=C["accent1"])

rect(s, 0.6, 3.75, 7, 0.04, "253E6A")

add_text(s, "AI 프로그래밍 팀 프로젝트  |  2026년 5월",
         0.6, 3.95, 9, 0.5, font_size=15, color=C["offWhite"])
add_text(s, "데이터 출처: 자동차관리정보시스템(국토교통부) · 서울특별시",
         0.6, 4.65, 10, 0.45, font_size=11, color=C["muted"])
add_text(s, "팀명 / 발표자 이름", 0.6, 5.2, 6, 0.45,
         font_size=12, italic=True, color=C["muted"])


# ════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "Agenda", "발표 순서")
page_num(s, 2, TOTAL)

agenda = [
    ("01", "주제 선정 배경",       "전기차 보급 vs 충전 인프라 불균형",      C["accent2"]),
    ("02", "데이터 수집",           "2개 공공데이터셋 소개",                  C["accent1"]),
    ("03", "데이터 전처리 & 지표",  "부족지수 · 충전편의지수 산출",           C["accent3"]),
    ("04", "서비스 주요 기능",      "지도 시각화 & 충전소 안내",              C["warn"]),
    ("05", "분석 결과",             "부족 지역 및 자치구별 현황",             C["pink"]),
    ("06", "기술 스택 & 시연",      "Python · Streamlit · Folium",           C["accent2"]),
    ("07", "결론 & Q&A",            "향후 발전 방향",                        C["accent3"]),
]

cols = [0.4, 7.1]
for i, (num, label, sub, color) in enumerate(agenda):
    col = i % 2
    row = i // 2
    x = cols[col]
    y = 1.55 + row * 1.6
    bw = 5.9

    rect(s, x, y, bw, 1.35, C["midBg"], "253E6A", 0.8)
    rect(s, x, y, 0.09, 1.35, color)
    add_text(s, num, x + 0.2, y + 0.1, 0.7, 0.55, font_size=22, bold=True, color=color)
    add_text(s, label, x + 1.0, y + 0.1, bw - 1.1, 0.58,
             font_size=16, bold=True, color=C["white"])
    add_text(s, sub, x + 1.0, y + 0.73, bw - 1.1, 0.48,
             font_size=11, color=C["muted"])


# ════════════════════════════════════════════════════════
# SLIDE 3: Background 1
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "왜 전기차 충전 인프라인가?", "주제 선정 배경 ①")
section_pill(s, "01  배경", C["accent2"])
page_num(s, 3, TOTAL)

# Stats
stats = [
    ("111,761", "서울시\n전기차 등록(대)", C["accent1"]),
    ("↑ 매년",  "전기차 등록 수\n지속 증가",  C["accent3"]),
    ("417곳",   "인프라 부족\n행정동 수", C["warn"]),
    ("데이터로", "인프라 현황\n정량 파악", C["pink"]),
]
for i, (val, label, color) in enumerate(stats):
    x = 0.4 + i * 3.2
    rect(s, x, 1.55, 2.9, 1.55, C["midBg"], "253E6A", 0.6)
    rect(s, x, 1.55, 2.9, 0.08, color)
    add_text(s, val, x + 0.1, 1.65, 2.7, 0.72,
             font_size=26, bold=True, color=color, align="center")
    add_text(s, label, x + 0.1, 2.35, 2.7, 0.65,
             font_size=10.5, color=C["offWhite"], align="center")

points = [
    "전기차 보급 급증 — 환경부 목표 2030년 친환경차 450만 대, 서울시 매년 가파른 증가세",
    '"충전 불안(Charge Anxiety)" → 충전소 부족이 전기차 구매 기피의 핵심 원인',
    "충전 인프라 확충 속도가 전기차 보급 속도를 따라가지 못하는 구조적 문제",
    "데이터 기반으로 어느 지역이 얼마나 부족한지 시각적·정량적으로 파악 필요",
]
for i, p in enumerate(points):
    rect(s, 0.4, 3.35 + i * 0.92, W - 0.8, 0.78, C["midBg"], "1E3560", 0.5)
    rect(s, 0.4, 3.35 + i * 0.92, 0.08, 0.78, C["accent2"])
    add_text(s, p, 0.65, 3.38 + i * 0.92, W - 1.1, 0.72,
             font_size=12, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 4: Problem Definition
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "문제 정의", "주제 선정 배경 ②")
section_pill(s, "01  배경", C["accent2"])
page_num(s, 4, TOTAL)

problems = [
    ("01", "데이터 기반 불균형 파악",
     "행정동 단위로 전기차 등록 수 vs 충전기 설치 수의 불균형을 정량적으로 측정하여 우선 개선 지역 식별",
     C["accent2"]),
    ("02", "직관적 지역 시각화",
     "인프라 부족 지역을 색상(빨강↔초록)으로 표현하여 비전문가도 한눈에 현황 파악 가능",
     C["accent1"]),
    ("03", "시민 활용 충전소 안내",
     "내 위치에서 가장 가까운 충전소를 즉시 검색 → 충전 불안 해소 및 실생활 편의 제공",
     C["accent3"]),
]
for i, (num, title, desc, color) in enumerate(problems):
    y = 1.6 + i * 1.7
    rect(s, 0.4, y, W - 0.8, 1.5, C["midBg"], "253E6A", 0.8)
    rect(s, 0.4, y, 0.09, 1.5, color)

    # Number circle
    circle = s.shapes.add_shape(9, I_(0.65), I_(y + 0.32), I_(0.68), I_(0.68))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(color)
    circle.line.fill.background()
    add_text(s, num, 0.65, y + 0.32, 0.68, 0.68,
             font_size=18, bold=True, color=C["darkBg"], align="center", valign="middle")

    add_text(s, title, 1.5, y + 0.12, W - 2.0, 0.55,
             font_size=17, bold=True, color=C["white"])
    add_text(s, desc, 1.5, y + 0.72, W - 2.0, 0.65,
             font_size=12, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 5: Data Collection
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "데이터 수집", "공공데이터 2종 활용")
section_pill(s, "02  데이터", C["accent1"])
page_num(s, 5, TOTAL)

half_w = (W - 1.0) / 2

# Dataset 1
rect(s, 0.4, 1.55, half_w, 5.6, C["midBg"], "253E6A", 0.7)
rect(s, 0.4, 1.55, half_w, 0.09, C["accent2"])
add_text(s, "📊  데이터셋 ①", 0.55, 1.65, half_w - 0.2, 0.5,
         font_size=14, bold=True, color=C["accent2"])
add_text(s, "서울시 자치구 읍면동별\n연료별 자동차 등록현황",
         0.55, 2.2, half_w - 0.2, 0.95,
         font_size=16, bold=True, color=C["white"])

d1 = [
    ("출처", "자동차관리정보시스템 (국토교통부)"),
    ("기준일", "2026년 4월"),
    ("형태", "Excel (.xlsx)  /  약 5,190행"),
    ("주요 컬럼", "시군구, 행정동, 연료, 등록대수"),
    ("전처리", "'전기' 연료 필터링 → 행정동별 합산"),
]
for i, (k, v) in enumerate(d1):
    y_pos = 3.25 + i * 0.75
    add_text(s, k, 0.55, y_pos, 1.5, 0.45,
             font_size=10, bold=True, color=C["accent2"])
    add_text(s, v, 2.05, y_pos, half_w - 1.85, 0.45,
             font_size=11, color=C["offWhite"])
    if i < len(d1) - 1:
        rect(s, 0.55, y_pos + 0.48, half_w - 0.3, 0.015, "253E6A")

# Dataset 2
x2 = 0.4 + half_w + 0.2
rect(s, x2, 1.55, half_w, 5.6, C["midBg"], "253E6A", 0.7)
rect(s, x2, 1.55, half_w, 0.09, C["accent1"])
add_text(s, "⚡  데이터셋 ②", x2 + 0.15, 1.65, half_w - 0.2, 0.5,
         font_size=14, bold=True, color=C["accent1"])
add_text(s, "서울시 전기차\n충전소 설치현황",
         x2 + 0.15, 2.2, half_w - 0.2, 0.95,
         font_size=16, bold=True, color=C["white"])

d2 = [
    ("출처", "서울특별시 제공"),
    ("기준일", "2026년 5월 12일"),
    ("형태", "Excel (.xlsx)  /  588개 충전소"),
    ("주요 컬럼", "설치장소, 주소, 급속·완속 충전기 수"),
    ("전처리", "주소에서 자치구·동명 추출 (정규표현식)"),
]
for i, (k, v) in enumerate(d2):
    y_pos = 3.25 + i * 0.75
    add_text(s, k, x2 + 0.15, y_pos, 1.5, 0.45,
             font_size=10, bold=True, color=C["accent1"])
    add_text(s, v, x2 + 1.65, y_pos, half_w - 1.85, 0.45,
             font_size=11, color=C["offWhite"])
    if i < len(d2) - 1:
        rect(s, x2 + 0.15, y_pos + 0.48, half_w - 0.3, 0.015, "253E6A")


# ════════════════════════════════════════════════════════
# SLIDE 6: Data Processing
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "데이터 전처리 과정", "두 데이터셋 결합 및 지표 산출")
section_pill(s, "03  전처리", C["accent3"])
page_num(s, 6, TOTAL)

steps = [
    ("1", "전기차\n필터링",    "연료='전기'\n→ 동별 합산",    C["accent2"]),
    ("2", "주소\n파싱",        "정규표현식으로\n구·동명 추출",  C["accent1"]),
    ("3", "데이터\nJOIN",      "구+동 기준\nOUTER JOIN",       C["accent3"]),
    ("4", "결측값\n처리",      "충전소 없는 동\n→ 0 처리",     C["warn"]),
    ("5", "지표\n산출",        "부족지수 &\n편의지수 계산",    C["pink"]),
]
step_w = (W - 0.8) / 5
for i, (num, label, desc, color) in enumerate(steps):
    x = 0.4 + i * step_w
    rect(s, x, 1.6, step_w - 0.15, 4.0, C["midBg"], "253E6A", 0.6)
    rect(s, x, 1.6, step_w - 0.15, 0.09, color)

    # Circle number
    cx = x + (step_w - 0.15) / 2 - 0.4
    circle = s.shapes.add_shape(9, I_(cx), I_(1.85), I_(0.8), I_(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(color)
    circle.line.fill.background()
    add_text(s, num, cx, 1.85, 0.8, 0.8,
             font_size=20, bold=True, color=C["darkBg"], align="center", valign="middle")

    add_text(s, label, x + 0.1, 2.82, step_w - 0.35, 0.75,
             font_size=13, bold=True, color=C["white"], align="center")
    add_text(s, desc, x + 0.1, 3.62, step_w - 0.35, 0.8,
             font_size=10.5, color=C["offWhite"], align="center")

    if i < len(steps) - 1:
        add_text(s, "›", x + step_w - 0.22, 2.85, 0.25, 0.5,
                 font_size=20, color=C["muted"], align="center")

# Formula box
rect(s, 0.4, 5.85, W - 0.8, 1.25, C["tagBg"], "253E6A", 0.8)
add_text(s, "부족지수 = 전기차 수 ÷ 충전기 수",
         0.6, 5.95, 5.5, 0.48, font_size=13, bold=True, color=C["warn"], align="center")
add_text(s, "|", 6.15, 5.95, 0.5, 1.05, font_size=28, color=C["muted"], align="center")
add_text(s, "충전편의지수 = (급속×3 + 완속) ÷ 전기차 수 × 100",
         6.7, 5.95, 6.2, 0.48, font_size=13, bold=True, color=C["accent1"], align="center")
add_text(s, "높을수록 인프라 부족",
         0.6, 6.45, 5.5, 0.4, font_size=10, color=C["muted"], italic=True, align="center")
add_text(s, "높을수록 충전 편리 (급속 가중치 3배)",
         6.7, 6.45, 6.2, 0.4, font_size=10, color=C["muted"], italic=True, align="center")


# ════════════════════════════════════════════════════════
# SLIDE 7: Key Metrics
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "핵심 분석 지표", "인프라 상태 정량화 기준")
section_pill(s, "03  전처리", C["accent3"])
page_num(s, 7, TOTAL)

half_w2 = (W - 1.0) / 2

# Shortage Index
rect(s, 0.4, 1.55, half_w2, 5.65, C["midBg"], "253E6A", 0.7)
rect(s, 0.4, 1.55, half_w2, 0.09, C["warn"])
add_text(s, "⚠️  부족지수 (Shortage Index)", 0.55, 1.65, half_w2 - 0.2, 0.5,
         font_size=15, bold=True, color=C["warn"])
add_text(s, "전기차 수 ÷ 충전기 수", 0.55, 2.2, half_w2 - 0.2, 0.45,
         font_size=13, italic=True, color=C["offWhite"])
add_text(s, "높을수록 인프라가 부족한 지역", 0.55, 2.65, half_w2 - 0.2, 0.4,
         font_size=11, color=C["muted"])

shortage_levels = [
    ("매우 부족", "> 20",   C["danger"]),
    ("부족",      "10 ~ 20", "F97316"),
    ("보통",      "5 ~ 10",  C["warn"]),
    ("충분",      "≤ 5",     C["accent3"]),
]
for i, (label, rng, color) in enumerate(shortage_levels):
    y_pos = 3.2 + i * 0.9
    rect(s, 0.55, y_pos, 0.22, 0.62, color)
    add_text(s, label, 0.88, y_pos, 2.2, 0.62,
             font_size=15, bold=True, color=color)
    add_text(s, f"({rng})", 3.1, y_pos, 2.8, 0.62,
             font_size=14, color=C["offWhite"])

# Convenience Index
x2 = 0.4 + half_w2 + 0.2
rect(s, x2, 1.55, half_w2, 5.65, C["midBg"], "253E6A", 0.7)
rect(s, x2, 1.55, half_w2, 0.09, C["accent1"])
add_text(s, "✅  충전편의지수 (Convenience Index)", x2 + 0.15, 1.65, half_w2 - 0.2, 0.5,
         font_size=14, bold=True, color=C["accent1"])
add_text(s, "(급속×3 + 완속) ÷ 전기차 수 × 100", x2 + 0.15, 2.2, half_w2 - 0.2, 0.45,
         font_size=12, italic=True, color=C["offWhite"])
add_text(s, "높을수록 충전이 편리한 지역 | 급속 가중치 3배", x2 + 0.15, 2.65, half_w2 - 0.2, 0.4,
         font_size=10.5, color=C["muted"])

conv_levels = [
    ("매우 좋음", "≥ 30",   C["accent3"]),
    ("좋음",      "≥ 15",   "6EE7B7"),
    ("보통",      "≥ 5",    C["warn"]),
    ("나쁨",      "0 ~ 5",  "F97316"),
    ("인프라 전무", "= 0",  C["danger"]),
]
for i, (label, rng, color) in enumerate(conv_levels):
    y_pos = 3.2 + i * 0.72
    rect(s, x2 + 0.15, y_pos, 0.22, 0.5, color)
    add_text(s, label, x2 + 0.48, y_pos, 2.5, 0.5,
             font_size=13.5, bold=True, color=color)
    add_text(s, f"({rng})", x2 + 3.0, y_pos, 2.5, 0.5,
             font_size=13, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 8: Service Features
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "서비스 주요 기능", "4가지 핵심 기능 소개")
section_pill(s, "04  기능", C["accent2"])
page_num(s, 8, TOTAL)

features = [
    ("🗺️ 01", "인프라 부족 지역 지도",
     ["행정동별 부족지수 색상 표현 (빨강 → 초록)",
      "자치구 choropleth 오버레이",
      "클릭 시 팝업: 전기차·충전기 수 확인"],
     C["accent2"], 0.4, 1.55),
    ("⚡ 02", "충전편의지수 지도",
     ["편의성 기반 지역 색상 표현",
      "급속·완속 가중치 반영한 색상",
      "4가지 모드 사이드바 토글 선택"],
     C["accent1"], 7.15, 1.55),
    ("📍 03", "가장 가까운 충전소 안내",
     ["주소 텍스트 입력 → 자동 위치 변환",
      "좌표 직접 입력도 지원",
      "가까운 충전소 Top 5 + 거리 표시"],
     C["accent3"], 0.4, 4.55),
    ("📊 04", "자치구·행정동 상세 테이블",
     ["자치구별 요약 (부족지수 순 정렬)",
      "상태 필터 (충전소 없음·부족 등)",
      "전체 행정동 데이터 테이블 제공"],
     C["warn"], 7.15, 4.55),
]

for icon_title, title, items, color, x, y in features:
    fw = W / 2 - 0.55
    rect(s, x, y, fw, 2.7, C["midBg"], "253E6A", 0.8)
    rect(s, x, y, fw, 0.09, color)
    add_text(s, f"{icon_title}  {title}", x + 0.15, y + 0.12, fw - 0.2, 0.58,
             font_size=15, bold=True, color=C["white"])
    for j, item in enumerate(items):
        add_text(s, f"›  {item}", x + 0.25, y + 0.78 + j * 0.6, fw - 0.4, 0.52,
                 font_size=12, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 9: Analysis Results
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "주요 분석 결과", "서울시 전기차 인프라 현황 (2026년 기준)")
section_pill(s, "05  결과", C["pink"])
page_num(s, 9, TOTAL)

kpis = [
    ("111,761", "전기차 등록(대)", C["accent1"]),
    ("1,471",   "충전기 설치(기)", C["accent3"]),
    ("588",     "충전소(개소)",    C["accent2"]),
    ("417곳",   "인프라 부족 동",  C["danger"]),
]
kw = (W - 0.8) / 4
for i, (val, label, color) in enumerate(kpis):
    x = 0.4 + i * kw
    rect(s, x, 1.55, kw - 0.15, 1.45, C["midBg"], "253E6A", 0.6)
    rect(s, x, 1.55, kw - 0.15, 0.09, color)
    add_text(s, val, x + 0.1, 1.65, kw - 0.35, 0.72,
             font_size=28, bold=True, color=color, align="center")
    add_text(s, label, x + 0.1, 2.35, kw - 0.35, 0.55,
             font_size=11, color=C["muted"], align="center")

findings = [
    ("🔴", "충전소 없는 동 TOP 3",
     "강남구 대치1동 (3,428대) · 강서구 가양1동 (2,800대) · 서초구 양재1동 (2,766대)",
     C["danger"]),
    ("📍", "부족지수 TOP 자치구",
     "은평구 (384.9) · 금천구 (295.4) · 강남구 (147.5) · 관악구 (137.7) · 강서구 (137.5)",
     C["warn"]),
    ("⚠️", "강남구 역설 (Paradox)",
     "전기차 최다 지역 (15,631대)이지만 충전기 106기에 불과 → 부족지수 147.5로 위험 수준",
     C["pink"]),
    ("📊", "충전소 편중 현상",
     "전체 충전소 588곳 중 강남·서초 집중 → 외곽 및 북부 지역 상대적 인프라 열악",
     C["accent2"]),
]

for i, (icon, head, body, color) in enumerate(findings):
    y_pos = 3.18 + i * 1.02
    rect(s, 0.4, y_pos, W - 0.8, 0.88, C["midBg"], "1E3560", 0.5)
    rect(s, 0.4, y_pos, 0.09, 0.88, color)
    add_text(s, f"{icon}  {head}", 0.6, y_pos + 0.06, 3.6, 0.4,
             font_size=12, bold=True, color=color)
    add_text(s, body, 4.3, y_pos + 0.06, W - 4.8, 0.72,
             font_size=11, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 10: Tech Stack
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "기술 구현 스택", "Python 기반 풀스택 데이터 시각화")
section_pill(s, "06  기술", C["accent2"])
page_num(s, 10, TOTAL)

stack = [
    ("🐍", "Python 3.11",          "주 개발 언어",                    C["accent2"]),
    ("🌐", "Streamlit",            "웹 애플리케이션 프레임워크",        C["accent1"]),
    ("🗺️", "Folium",              "인터랙티브 지도 시각화",            C["accent3"]),
    ("📦", "Pandas / OpenPyXL",    "데이터 처리 & Excel 파싱",         C["warn"]),
    ("📍", "Geopy (Nominatim)",    "주소 → 좌표 변환 (Geocoding)",     C["pink"]),
    ("📐", "Haversine Formula",    "두 지점 간 거리 계산 (km)",        C["accent2"]),
]

sw = (W - 0.8) / 3
for i, (icon, name, role, color) in enumerate(stack):
    col = i % 3
    row = i // 3
    x = 0.4 + col * (sw + 0.1)
    y = 1.6 + row * 2.0
    rect(s, x, y, sw - 0.05, 1.7, C["midBg"], "253E6A", 0.7)

    circle = s.shapes.add_shape(9, I_(x + 0.25), I_(y + 0.3), I_(0.75), I_(0.75))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(color)
    circle.line.fill.background()
    add_text(s, icon, x + 0.25, y + 0.3, 0.75, 0.75,
             font_size=20, align="center", valign="middle", color=C["white"])

    add_text(s, name, x + 1.15, y + 0.25, sw - 1.3, 0.55,
             font_size=14, bold=True, color=C["white"])
    add_text(s, role, x + 1.15, y + 0.82, sw - 1.3, 0.7,
             font_size=11, color=C["muted"])


# ════════════════════════════════════════════════════════
# SLIDE 11: Demo
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["midBg"])

rect(s, 0, 0, W, 1.3, C["darkBg"])
rect(s, 0, 1.3, W, 0.07, C["accent3"])
add_text(s, "🖥️  서비스 시연 (Live Demo)", 0.5, 0.12, W - 1, 1.05,
         font_size=32, bold=True, color=C["white"])
page_num(s, 11, TOTAL)

screens = [
    ("부족지수 지도\n색상 시각화", 0.3, 4.0),
    ("편의지수 지도\n4모드 토글",  4.55, 4.0),
    ("충전소 검색\nTop 5 안내",   8.8, 4.0),
]
sw2 = 4.0
for title, x, w in screens:
    rect(s, x, 1.58, w, 4.2, "0A1525", C["accent2"], 1.2)
    add_text(s, "[ 화면 캡처 이미지 ]", x + 0.3, 2.8, w - 0.6, 0.7,
             font_size=13, color="3A5580", align="center", italic=True)
    rect(s, x, 5.55, w, 0.9, C["cardBg"], "253E6A", 0.6)
    add_text(s, title, x + 0.1, 5.58, w - 0.2, 0.82,
             font_size=12, bold=True, color=C["accent1"], align="center")

add_text(s, "실행 방법:  streamlit run app.py   →   브라우저에서 http://localhost:8501 접속",
         0.5, 6.6, W - 1, 0.55,
         font_size=12, color=C["accent3"], align="center", italic=True,
         font_name="Consolas")


# ════════════════════════════════════════════════════════
# SLIDE 12: Conclusion & Future
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
add_title_bar(s, "결론 및 향후 발전 방향", "")
section_pill(s, "07  결론", C["accent3"])
page_num(s, 12, TOTAL)

half_w3 = (W - 1.0) / 2

rect(s, 0.4, 1.55, half_w3, 5.65, C["midBg"], "253E6A", 0.7)
rect(s, 0.4, 1.55, half_w3, 0.09, C["accent3"])
add_text(s, "✅  결론", 0.55, 1.65, half_w3 - 0.2, 0.5,
         font_size=16, bold=True, color=C["accent3"])

conclusions = [
    "공공데이터 기반으로 충전 인프라 불균형 정량적 확인",
    "강남·강서 고밀도 주거지역 충전 확충 필요성 도출",
    "부족지수·편의지수로 지역별 인프라 수준 비교 가능",
    "시민이 활용 가능한 충전소 안내 서비스 구현",
    "Python + 오픈소스로 실용적 데이터 서비스 완성",
]
for i, c in enumerate(conclusions):
    add_text(s, f"✓  {c}", 0.55, 2.3 + i * 0.82, half_w3 - 0.25, 0.72,
             font_size=12, color=C["offWhite"])

x2 = 0.4 + half_w3 + 0.2
rect(s, x2, 1.55, half_w3, 5.65, C["midBg"], "253E6A", 0.7)
rect(s, x2, 1.55, half_w3, 0.09, C["accent2"])
add_text(s, "🚀  향후 발전 방향", x2 + 0.15, 1.65, half_w3 - 0.2, 0.5,
         font_size=16, bold=True, color=C["accent2"])

futures = [
    ("📡", "실시간 충전기 가용 API 연동 (환경부 공공API)"),
    ("📍", "충전소 정확 좌표 확보 → Kakao/Naver Geocoding"),
    ("🧭", "이동 경로 기반 충전소 추천 (네비게이션 연동)"),
    ("🔮", "AI 기반 미래 충전 수요 예측 모델 적용"),
    ("🇰🇷", "서울시 외 전국 광역시·도 확장"),
    ("📱", "모바일 앱 버전 개발 (React Native 등)"),
]
for i, (icon, text) in enumerate(futures):
    add_text(s, f"{icon}  {text}", x2 + 0.15, 2.3 + i * 0.82, half_w3 - 0.3, 0.72,
             font_size=12, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 13: Q&A
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["darkBg"])

circle1 = s.shapes.add_shape(9, I_(-1.5), I_(4), I_(6), I_(6))
circle1.fill.solid()
circle1.fill.fore_color.rgb = rgb("001A40")
circle1.line.fill.background()

circle2 = s.shapes.add_shape(9, I_(10), I_(-1), I_(4.5), I_(4.5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = rgb("001A3A")
circle2.line.fill.background()

rect(s, 0, 0, 0.35, H, C["accent1"])

add_text(s, "Q & A", 1.0, 1.3, W - 2.0, 2.0,
         font_size=72, bold=True, color=C["white"], align="center")

rect(s, 3.0, 3.9, 7.33, 0.05, "253E6A")

add_text(s, "감사합니다", 1.0, 4.1, W - 2.0, 0.85,
         font_size=32, color=C["offWhite"], align="center")
add_text(s, "서울시 전기차 충전 인프라 부족 지역 분석 및 시각화 서비스",
         1.0, 5.1, W - 2.0, 0.55,
         font_size=14, color=C["muted"], align="center", italic=True)

rect(s, 4.0, 6.0, 5.33, 0.9, C["tagBg"], "253E6A", 0.8)
add_text(s, "팀명 / 발표자 이름 입력", 4.0, 6.0, 5.33, 0.9,
         font_size=13, color=C["muted"], align="center", valign="middle", italic=True)

page_num(s, 13, TOTAL)


# ─── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT_FILE)
print(f"✅ PPT 생성 완료: {OUT_FILE}")
