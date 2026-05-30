"""
주제 제안 발표용 PPT - 스토리텔링 구조
"전기차를 사려고 했더니..."
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches as I_, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_FILE = "서울시_EV_인프라분석_주제제안.pptx"

# ─── Color palette ────────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C = {
    "darkBg":    "060E1F",
    "midBg":     "0D1B35",
    "cardBg":    "1A2B4A",
    "contentBg": "EEF2FF",
    "accent1":   "00C8FF",
    "accent2":   "3B82F6",
    "accent3":   "10B981",
    "warn":      "F59E0B",
    "danger":    "EF4444",
    "white":     "FFFFFF",
    "offWhite":  "D0DCF0",
    "muted":     "7A94BF",
    "dark":      "1E2B42",
    "tagBg":     "142038",
    "pink":      "F472B6",
    "orange":    "F97316",
}

prs = Presentation()
prs.slide_width  = I_(13.33)
prs.slide_height = I_(7.5)
W, H = 13.33, 7.5
BLANK = prs.slide_layouts[6]
TOTAL = 12


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill_hex, line_hex=None, lw=0):
    s = slide.shapes.add_shape(1, I_(x), I_(y), I_(w), I_(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex and lw > 0:
        s.line.color.rgb = rgb(line_hex)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def oval(slide, x, y, w, h, fill_hex, line_hex=None, lw=0):
    s = slide.shapes.add_shape(9, I_(x), I_(y), I_(w), I_(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex and lw > 0:
        s.line.color.rgb = rgb(line_hex)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h,
        fs=12, bold=False, italic=False,
        color="FFFFFF", align="left", valign="middle",
        font="맑은 고딕", wrap=True):
    tb = slide.shapes.add_textbox(I_(x), I_(y), I_(w), I_(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = {
        "middle": 3, "top": 1, "bottom": 4
    }.get(valign, 3)
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = {"center": PP_ALIGN.CENTER,
                   "right":  PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    run.font.name = font
    return tb


def header_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, W, 1.35, C["midBg"])
    rect(slide, 0, 1.35, W, 0.065, C["accent1"])
    txt(slide, title, 0.55, 0.1, W - 1.1, 0.88,
        fs=30, bold=True, color=C["white"])
    if subtitle:
        txt(slide, subtitle, 0.55, 0.92, W - 1.1, 0.38,
            fs=12, color=C["offWhite"])


def step_badge(slide, label, color):
    rect(slide, 0.5, 0.28, 2.2, 0.42, C["cardBg"], color, 1.0)
    txt(slide, label, 0.5, 0.28, 2.2, 0.42,
        fs=9, bold=True, color=color, align="center")


def page_num(slide, n):
    txt(slide, f"{n} / {TOTAL}", W - 1.3, H - 0.45, 0.95, 0.32,
        fs=9, color=C["muted"], align="right")


def quote_box(slide, text, x, y, w, h, color):
    """인용문 스타일 강조 박스"""
    rect(slide, x, y, 0.1, h, color)
    rect(slide, x, y, w, h, C["cardBg"], "253E6A", 0.5)
    txt(slide, text, x + 0.25, y, w - 0.3, h,
        fs=14, italic=True, color=C["offWhite"], valign="middle")


def card(slide, x, y, w, h, accent_color, title, body_lines, title_fs=14):
    rect(slide, x, y, w, h, C["midBg"], "253E6A", 0.7)
    rect(slide, x, y, w, 0.09, accent_color)
    txt(slide, title, x + 0.2, y + 0.12, w - 0.3, 0.5,
        fs=title_fs, bold=True, color=C["white"])
    for i, line in enumerate(body_lines):
        txt(slide, line, x + 0.2, y + 0.68 + i * 0.52, w - 0.3, 0.48,
            fs=11.5, color=C["offWhite"])


def stat_box(slide, x, y, w, h, val, label, color):
    rect(slide, x, y, w, h, C["midBg"], "253E6A", 0.6)
    rect(slide, x, y, w, 0.09, color)
    txt(slide, val,   x + 0.1, y + 0.1,  w - 0.2, 0.72,
        fs=30, bold=True, color=color, align="center")
    txt(slide, label, x + 0.1, y + 0.82, w - 0.2, 0.52,
        fs=11, color=C["muted"], align="center")


# ════════════════════════════════════════════════════════
# SLIDE 1: 표지
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["darkBg"])

# 배경 장식
oval(s, 8.0, -1.2, 6.5, 6.5, "00152E")
oval(s, -1.0, 4.5, 5.0, 5.0, "001530")
rect(s, 0, 0, 0.35, H, C["accent1"])

# 태그 라벨
rect(s, 0.6, 0.55, 3.2, 0.42, C["tagBg"], C["accent2"], 1.0)
txt(s, "AI 프로그래밍  |  주제 제안 발표", 0.6, 0.55, 3.2, 0.42,
    fs=9, bold=True, color=C["accent2"], align="center")

# 메인 타이틀
txt(s, "전기차를 사려고 했더니...", 0.6, 1.3, W - 1.0, 1.0,
    fs=14, italic=True, color=C["muted"])
txt(s, "서울시 전기차 충전 인프라", 0.6, 2.1, W - 1.0, 1.1,
    fs=44, bold=True, color=C["white"])
txt(s, "부족 지역 분석 및 시각화 서비스", 0.6, 3.15, W - 1.0, 0.85,
    fs=32, bold=True, color=C["accent1"])

rect(s, 0.6, 4.18, 8.5, 0.05, "253E6A")

txt(s, "공공데이터를 활용한 인프라 불균형 분석 & 충전소 안내 웹 서비스 제안",
    0.6, 4.35, W - 1.0, 0.5, fs=13, color=C["offWhite"])
txt(s, "2026년 5월  |  팀명 / 발표자 이름",
    0.6, 4.98, W - 1.0, 0.42, fs=11, color=C["muted"])


# ════════════════════════════════════════════════════════
# SLIDE 2: 이야기의 시작 — 전기차 구매 고민
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "이야기의 시작", "프로젝트 주제를 선정하게 된 계기")
step_badge(s, "STEP 01  동기", C["accent2"])
page_num(s, 2)

quote_box(s,
    '"요즘 전기차가 많아졌던데... 나도 한 번 바꿔볼까?"',
    0.55, 1.55, W - 1.1, 0.88, C["accent1"])

reasons = [
    ("💰", "경제성",  "유류비 절감, 세제 혜택, 보조금 지원 등 경제적 매력 증가"),
    ("🌱", "환경",    "탄소 배출 저감 — 환경부 2030 무공해차 450만 대 보급 목표"),
    ("🚗", "트렌드",  "주변에서 전기차를 타는 사람이 눈에 띄게 늘어난 체감"),
]
for i, (icon, label, desc) in enumerate(reasons):
    y = 2.65 + i * 1.46
    rect(s, 0.55, y, W - 1.1, 1.26, C["midBg"], "253E6A", 0.7)
    oval(s, 0.75, y + 0.28, 0.72, 0.72, C["accent2"])
    txt(s, icon, 0.75, y + 0.28, 0.72, 0.72,
        fs=22, align="center", valign="middle", color=C["white"])
    txt(s, label, 1.62, y + 0.1, 2.2, 0.5,
        fs=16, bold=True, color=C["accent2"])
    txt(s, desc, 1.62, y + 0.62, W - 2.4, 0.52,
        fs=12, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 3: 그런데 걱정이 생겼다 — 충전 불안
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "그런데 걱정이 생겼다", "막상 사려니 이런 궁금증이 — 실제 설문에서도 같은 이유가 TOP을 차지했다")
step_badge(s, "STEP 02  문제인식", C["warn"])
page_num(s, 3)

questions = [
    ("❓", "주변에 전기차 타는 사람이 얼마나 될까?",
     "우리 동네·구에 전기차가 얼마나 보급되어 있을까?\n→ 서울시만 111,761대 — 생각보다 훨씬 많다",
     C["accent2"]),
    ("❓", "충전소가 가까이 있기는 한 걸까?",
     "집 근처에 충전소가 없으면 매일 어디서 충전해야 하지?\n→ 소비자 38.6%가 충전 인프라 부족을 1순위 불만으로 꼽음",
     C["warn"]),
    ("❓", "충전하려고 줄 서서 기다리면 어떡하지?",
     "전기차 늘수록 충전기 대기·주행거리 불안 가중\n→ 충전 시간 우려 49%  /  주행거리 불안 26.4%",
     C["danger"]),
    ("❓", "배터리 화재·급발진이 걱정된다",
     '"충전 불안 · 안전 불안" — 전기차 구매 기피의 핵심\n→ 배터리 안전성 우려 50%  (2025년 신규 1위 등극)',
     C["pink"]),
]
for i, (icon, q, sub, color) in enumerate(questions):
    col = i % 2
    row = i // 2
    x = 0.55 + col * 6.45
    y = 1.6 + row * 2.52
    w = 6.1
    rect(s, x, y, w, 2.22, C["midBg"], "253E6A", 0.7)
    rect(s, x, y, w, 0.09, color)
    txt(s, icon, x + 0.2, y + 0.12, 0.65, 0.58,
        fs=22, color=color)
    txt(s, q, x + 1.0, y + 0.12, w - 1.15, 0.62,
        fs=13, bold=True, color=C["white"])
    txt(s, sub, x + 0.2, y + 0.82, w - 0.4, 1.22,
        fs=11, italic=True, color=C["offWhite"])

# 출처 각주
txt(s,
    "출처: 한국환경공단 전기차 이용 실태조사 2024 (n=1,000) | 딜로이트 글로벌 자동차 소비자 설문 2025 (n=31,000+, 30개국) | EV트렌드코리아 2025 (n=8,072)",
    0.55, 6.62, W - 1.1, 0.36,
    fs=8, italic=True, color=C["muted"])


# ════════════════════════════════════════════════════════
# SLIDE 4: 공공데이터로 알아보기
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "공공데이터로 직접 확인해봤다", "두 가지 공공데이터셋 발굴")
step_badge(s, "STEP 03  데이터", C["accent1"])
page_num(s, 4)

quote_box(s,
    '"궁금하면 데이터로 확인해보자 — 공공데이터포털에서 찾아봤다"',
    0.55, 1.55, W - 1.1, 0.82, C["accent3"])

# 3개 데이터셋 카드
cw = (W - 1.3) / 3
datasets = [
    (C["accent2"], "📊  데이터셋 ①",
     "서울시 읍면동별 연료별\n자동차 등록현황 (행정동)",
     [
         "출처 : 자동차관리정보시스템 (국토교통부)",
         "기준일 : 2026년 4월",
         "5,190여 행  /  Excel (.xlsx) 형태",
         "구·동별 연료 종류와 등록 대수 포함",
         "→ '전기' 필터 후 동별 합산 처리",
     ]),
    (C["warn"], "📋  데이터셋 ②",
     "서울시 전기차\n충전소 설치현황 (Excel)",
     [
         "출처 : 서울특별시 제공",
         "기준일 : 2026년 5월 12일",
         "588개 충전소  /  Excel (.xlsx) 형태",
         "설치장소, 주소, 급속·완속 충전기 수 포함",
         "→ 주소 파싱으로 구·동 위치 추출",
     ]),
    (C["accent1"], "⚡  데이터셋 ③",
     "한국환경공단\n전기차 충전소 API",
     [
         "출처 : 공공데이터포털 (한국환경공단)",
         "기준일 : 2026년 5월 (실시간 갱신)",
         "70,168기  /  REST JSON API",
         "충전소명, 주소, GPS 좌표, 충전기 타입 포함",
         "→ GPS 기반 행정동 폴리곤 매핑으로 정밀 분석",
     ]),
]
for i, (color, header, title, items) in enumerate(datasets):
    x = 0.55 + i * (cw + 0.1)
    rect(s, x, 2.58, cw, 4.55, C["midBg"], "253E6A", 0.7)
    rect(s, x, 2.58, cw, 0.1, color)
    txt(s, header, x + 0.17, 2.7, cw - 0.25, 0.52,
        fs=13, bold=True, color=color)
    txt(s, title, x + 0.17, 3.28, cw - 0.25, 0.88,
        fs=13.5, bold=True, color=C["white"])
    for j, line in enumerate(items):
        txt(s, f"  ▪  {line}", x + 0.17, 4.28 + j * 0.48, cw - 0.25, 0.44,
            fs=10.5, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 5: 서울시 전기차 현황 — 예상보다 많다
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "서울시 전기차, 생각보다 훨씬 많았다", "데이터로 확인한 서울시 전기차 현황 (2026.04)")
step_badge(s, "STEP 04  현황 파악", C["accent3"])
page_num(s, 5)

# 핵심 수치
kw = (W - 1.1) / 4
for i, (val, label, color) in enumerate([
    ("111,761대",  "서울시 전기차\n총 등록 대수",    C["accent1"]),
    ("427개 동",   "전기차가 있는\n행정동 수",       C["accent2"]),
    ("15,631대",   "강남구 최다\n(서울시 1위)",      C["accent3"]),
    ("매년 ↑",     "전기차 등록 수\n지속 증가 추세", C["warn"]),
]):
    x = 0.55 + i * (kw + 0.03)
    rect(s, x, 1.58, kw, 1.55, C["midBg"], "253E6A", 0.6)
    rect(s, x, 1.58, kw, 0.09, color)
    txt(s, val, x + 0.1, 1.68, kw - 0.2, 0.72,
        fs=22, bold=True, color=color, align="center")
    txt(s, label, x + 0.1, 2.38, kw - 0.2, 0.62,
        fs=11, color=C["muted"], align="center")

# 스토리 텍스트
quote_box(s,
    '"서울에만 이미 11만 명 이상이 전기차를 타고 있구나.\n  그럼 이 사람들은 다 어디서 충전하고 있는 걸까?"',
    0.55, 3.38, W - 1.1, 1.1, C["accent1"])

# TOP 5 구 리스트
txt(s, "자치구별 전기차 등록 TOP 5", 0.55, 4.65, 5.0, 0.42,
    fs=12, bold=True, color=C["white"])
top5 = [
    ("1위", "강남구",  "15,631대", C["accent1"]),
    ("2위", "송파구",  "8,900대",  C["accent2"]),
    ("3위", "강서구",  "8,251대",  C["accent3"]),
    ("4위", "서초구",  "7,500대",  C["warn"]),
    ("5위", "노원구",  "6,200대",  C["muted"]),
]
for i, (rank, gu, cnt, color) in enumerate(top5):
    x = 0.55 + i * 2.52
    rect(s, x, 5.15, 2.35, 0.88, C["cardBg"], "253E6A", 0.5)
    txt(s, rank, x + 0.1, 5.18, 0.7, 0.38, fs=10, bold=True, color=color)
    txt(s, gu,   x + 0.1, 5.52, 1.4, 0.42, fs=13, bold=True, color=C["white"])
    txt(s, cnt,  x + 1.45, 5.18, 0.8, 0.72, fs=11, color=color, align="right")


# ════════════════════════════════════════════════════════
# SLIDE 6: 충전 인프라 현황 — 충격적인 현실
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "그런데 충전 인프라는?", "한국환경공단 API 전수 데이터로 확인한 서울시 충전소 현황")
step_badge(s, "STEP 04  현황 파악", C["accent3"])
page_num(s, 6)

# 대조 강조 수치
txt(s, "전기차  111,761 대", 0.55, 1.58, 5.9, 0.88,
    fs=30, bold=True, color=C["accent1"])
txt(s, "충전기   70,168 기", 7.0, 1.58, 5.9, 0.88,
    fs=30, bold=True, color=C["accent3"])
rect(s, 0.55, 2.55, 12.23, 0.055, "253E6A")

# 비율 강조
rect(s, 0.55, 2.75, 12.23, 1.52, C["midBg"], "253E6A", 0.7)
txt(s, "전기차 1.59대당 충전기 1기", 0.75, 2.88, 6.0, 0.55,
    fs=16, bold=True, color=C["accent3"])
txt(s, "전체 수는 충분 — 그러나 지역별 분포가 문제", 0.75, 3.38, 7.0, 0.55,
    fs=13, italic=True, color=C["warn"])
txt(s, "급속 5,760기  /  완속 64,408기  /  충전소 12,549곳", 7.8, 2.98, 5.0, 0.75,
    fs=12, color=C["muted"], align="center")

# 소비자 설문 인용 (3개 수치 배지)
survey_stats = [
    ("38.6%", "충전 인프라 부족\n불만 1순위",     C["danger"],  "한국환경공단 2024"),
    ("49%",   "충전 시간 오래\n걸림 우려",         C["warn"],    "딜로이트 GACS 2025"),
    ("50%",   "배터리 안전성\n우려",               C["orange"],  "딜로이트 GACS 2025"),
    ("14%",   "BEV 구매 선호도\n(하이브리드 31%)", C["pink"],    "딜로이트 GACS 2025"),
]
sw = (W - 1.1) / 4
for i, (pct, label, color, src) in enumerate(survey_stats):
    x = 0.55 + i * (sw + 0.02)
    rect(s, x, 4.42, sw, 1.48, C["midBg"], "253E6A", 0.6)
    rect(s, x, 4.42, sw, 0.09, color)
    txt(s, pct,   x + 0.1, 4.52, sw - 0.2, 0.58,
        fs=26, bold=True, color=color, align="center")
    txt(s, label, x + 0.1, 5.1,  sw - 0.2, 0.52,
        fs=10.5, color=C["white"], align="center")
    txt(s, src,   x + 0.1, 5.65, sw - 0.2, 0.2,
        fs=8.5, color=C["muted"], align="center")

# 추가 현황 카드
findings = [
    ("44곳",     "충전소 전무\n행정동 수",              C["danger"], "GPS 기반 정밀 분석 결과 — 실제 충전소 없는 행정동 44곳"),
    ("42%",      "상위 5개 구\n전기차 집중도",          C["warn"],   "강남·서초·송파·강서·강동 5개 구에 전기차 42% 집중"),
    ("강남구",   "부족지수\n서울 1위 (2.8)",            C["orange"], "강남구: 전기차 15,631대 / 충전기 5,665기 — 부족지수 2.8"),
    ("+4,755기", "강남구 인프라\n정상화 필요 충전기",   C["pink"],   "강남구 부족지수 1.5 수준 달성에 추가로 필요한 충전기 수"),
]
fw = (W - 1.1) / 4
for i, (val, label, color, desc) in enumerate(findings):
    x = 0.55 + i * (fw + 0.02)
    rect(s, x, 6.05, fw, 1.1, C["midBg"], "253E6A", 0.6)
    rect(s, x, 6.05, fw, 0.09, color)
    txt(s, val,   x + 0.1, 6.14, fw - 0.2, 0.4,
        fs=16, bold=True, color=color, align="center")
    txt(s, label, x + 0.1, 6.54, fw - 0.2, 0.5,
        fs=9.5, bold=True, color=C["white"], align="center")

# 출처 각주
txt(s,
    "출처: 한국환경공단 전기차 이용 실태조사 2024 (n=1,000) | 딜로이트 GACS 2025 (30개국 31,000명+, 한국 906명, 2025.01 발표) | EV트렌드코리아 2025 (n=8,072)",
    0.55, 7.18, W - 1.1, 0.26,
    fs=7.5, italic=True, color=C["muted"])


# ════════════════════════════════════════════════════════
# SLIDE 7: 지역별 불균형 — 이건 심각하다
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "지역별 불균형이 심각하다", "자치구별 전기차 vs 충전기 비교")
step_badge(s, "STEP 05  분석", C["danger"])
page_num(s, 7)

quote_box(s,
    '"절대 수는 개선됐지만 지역별 편차는 여전하다\n  — 강남구가 전기차 최다임에도 충전기 대비 부족 1위"',
    0.55, 1.55, W - 1.1, 0.98, C["danger"])

# 상위 부족 구 비교 바 차트 (텍스트 기반)
txt(s, "부족지수 상위 자치구  (= 전기차 수 ÷ 충전기 수, 높을수록 부족)  ※ API 전수 데이터 기준",
    0.55, 2.72, W - 1.1, 0.45, fs=12, bold=True, color=C["dark"])

gu_data = [
    ("강남구",  2.8, 15631, 5665,  C["danger"]),
    ("강동구",  2.1, 5895,  2788,  C["danger"]),
    ("관악구",  2.1, 3166,  1510,  C["danger"]),
    ("양천구",  2.0, 4131,  2097,  C["orange"]),
    ("광진구",  1.9, 2595,  1400,  C["orange"]),
    ("서초구",  1.9, 8870,  4654,  C["warn"]),
    ("강서구",  1.8, 8251,  4509,  C["warn"]),
]
max_idx = 2.8
bar_max_w = 7.5
for i, (gu, idx, ev, ch, color) in enumerate(gu_data):
    y = 3.32 + i * 0.58
    # 배경 줄 (가독성)
    rect(s, 0.55, y + 0.02, W - 1.1, 0.44, "E8ECF5")
    txt(s, gu, 0.58, y, 1.3, 0.5, fs=12, bold=True, color=C["dark"])
    bar_w = (idx / max_idx) * bar_max_w
    rect(s, 1.95, y + 0.06, bar_w, 0.34, color)
    txt(s, f"{idx}  (전기차 {ev:,}대 / 충전기 {ch:,}기)",
        2.08 + bar_w, y + 0.08, 10.0, 0.34,
        fs=10.5, color=C["dark"])


# ════════════════════════════════════════════════════════
# SLIDE 8: 기존 서비스 한계 + 문제 정의
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "기존 서비스로는 알 수 없었다", "충전소 안내는 있지만, 인프라 불균형 분석은 없다")
step_badge(s, "STEP 05  문제 정의", C["pink"])
page_num(s, 8)

# 기존 서비스 비교 테이블
rect(s, 0.55, 1.55, W - 1.1, 0.52, C["tagBg"], "253E6A", 0.7)
for j, (label, x) in enumerate([
    ("서비스",       0.75),
    ("충전소 위치",   3.6),
    ("실시간 가용",   6.05),
    ("지역별 통계",   8.5),
    ("전기차 밀도 대비\n부족 분석",  10.5),
]):
    txt(s, label, x, 1.58, 2.1, 0.46, fs=10, bold=True, color=C["muted"], align="center")

services = [
    ("카카오맵 / 네이버맵",   "✅", "✅", "❌", "❌", C["accent2"]),
    ("환경부 EV Where",       "✅", "✅", "△",  "❌", C["accent1"]),
    ("차지비·에버온 등 앱",   "✅", "✅", "❌", "❌", C["accent3"]),
    ("차지인포 (ksga.org)",   "✅", "❌", "✅", "❌", C["warn"]),
    ("PlugShare DataTool",    "✅", "△",  "✅", "△",  C["muted"]),
]
for i, (name, c1, c2, c3, c4, color) in enumerate(services):
    y = 2.2 + i * 0.56
    bg = C["midBg"] if i % 2 == 0 else C["cardBg"]
    rect(s, 0.55, y, W - 1.1, 0.53, bg, "1E3560", 0.4)
    rect(s, 0.55, y, 0.06, 0.53, color)
    txt(s, name, 0.72, y + 0.05, 2.75, 0.43, fs=11, color=C["white"])
    for val, xpos in [(c1, 3.6), (c2, 6.05), (c3, 8.5), (c4, 10.5)]:
        vc = C["accent3"] if val == "✅" else (C["warn"] if val == "△" else C["danger"])
        txt(s, val, xpos, y + 0.05, 2.1, 0.43, fs=14, color=vc, align="center")

txt(s, "△ = 일부 제공 또는 유료 B2B 서비스",
    0.55, 5.02, 6.0, 0.35, fs=9, italic=True, color=C["muted"])

# 핵심 GAP 강조
rect(s, 0.55, 5.48, W - 1.1, 1.65, C["midBg"], C["danger"], 1.2)
rect(s, 0.55, 5.48, 0.1, 1.65, C["danger"])
txt(s, "채워지지 않은 공백", 0.78, 5.55, 4.0, 0.48,
    fs=13, bold=True, color=C["danger"])
txt(s,
    '"전기차 등록 수 대비 충전기가 얼마나 부족한가" — 이 질문에 답하는\n'
    '행정동 단위 정량 분석 + 시각화 서비스는 국내에 존재하지 않는다',
    0.78, 6.06, W - 1.4, 0.98,
    fs=12.5, color=C["offWhite"])


# ════════════════════════════════════════════════════════
# SLIDE 9: 우리가 제안하는 것
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "우리가 만들 것", "데이터 기반 전기차 충전 인프라 시각화 서비스")
step_badge(s, "STEP 06  서비스 제안", C["accent3"])
page_num(s, 9)

# 서비스 한 줄 정의
rect(s, 0.55, 1.55, W - 1.1, 0.9, C["tagBg"], C["accent3"], 1.2)
txt(s, "서울시 전기차 충전 인프라 부족 지역을 분석하고,\n시민이 직접 활용할 수 있는 웹 기반 시각화 서비스",
    0.75, 1.58, W - 1.3, 0.84, fs=14, bold=True, color=C["white"], align="center")

features = [
    ("🗺️", "인프라 지도 시각화",
     "행정동별 부족지수를 색상으로 표현\n빨강(부족) ↔ 초록(충분) 직관적 표시",
     C["accent2"]),
    ("📊", "충전편의지수 제공",
     "급속·완속 가중치 반영한 편의지수 산출\n지역별 충전 편의성 비교 가능",
     C["accent1"]),
    ("📍", "가까운 충전소 안내",
     "내 주소 입력 → 가까운 충전소 Top 5\n거리·충전기 수 즉시 확인",
     C["accent3"]),
    ("📋", "상세 현황 테이블",
     "자치구·행정동별 전기차·충전기 수\n부족도 현황 데이터 제공",
     C["warn"]),
]
fw2 = (W - 1.1) / 4
for i, (icon, title, desc, color) in enumerate(features):
    x = 0.55 + i * (fw2 + 0.02)
    rect(s, x, 2.68, fw2, 4.45, C["midBg"], "253E6A", 0.7)
    rect(s, x, 2.68, fw2, 0.09, color)
    oval(s, x + fw2/2 - 0.45, 2.88, 0.9, 0.9, color)
    txt(s, icon, x + fw2/2 - 0.45, 2.88, 0.9, 0.9,
        fs=24, align="center", valign="middle", color=C["white"])
    txt(s, title, x + 0.15, 3.9,  fw2 - 0.3, 0.62,
        fs=13, bold=True, color=C["white"], align="center")
    txt(s, desc,  x + 0.15, 4.58, fw2 - 0.3, 1.4,
        fs=11, color=C["offWhite"], align="center")


# ════════════════════════════════════════════════════════
# SLIDE 10: 데이터 & 기술 계획
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "어떻게 만들 것인가", "데이터 처리 흐름 & 기술 스택")
step_badge(s, "STEP 07  구현 계획", C["accent2"])
page_num(s, 10)

# 흐름도 (5단계)
steps = [
    ("📥", "데이터\n수집",     "공공데이터\n2종 Excel",     C["accent2"]),
    ("🔧", "데이터\n전처리",   "연료 필터링\n주소 파싱",    C["accent1"]),
    ("📐", "지수 산출",       "부족지수\n편의지수",         C["accent3"]),
    ("🗺️", "지도\n시각화",    "Folium\nChoropleth",        C["warn"]),
    ("🌐", "웹 서비스",       "Streamlit\n배포",           C["pink"]),
]
sw = (W - 1.1) / 5
for i, (icon, title, desc, color) in enumerate(steps):
    x = 0.55 + i * sw
    rect(s, x, 1.55, sw - 0.12, 3.1, C["midBg"], "253E6A", 0.6)
    rect(s, x, 1.55, sw - 0.12, 0.09, color)
    oval(s, x + (sw - 0.12)/2 - 0.42, 1.75, 0.85, 0.85, color)
    txt(s, icon, x + (sw - 0.12)/2 - 0.42, 1.75, 0.85, 0.85,
        fs=22, align="center", valign="middle", color=C["white"])
    txt(s, title, x + 0.12, 2.72, sw - 0.35, 0.7,
        fs=12.5, bold=True, color=C["white"], align="center")
    txt(s, desc,  x + 0.12, 3.45, sw - 0.35, 0.6,
        fs=11, color=C["muted"], align="center")
    if i < len(steps) - 1:
        txt(s, "→", x + sw - 0.18, 2.82, 0.25, 0.42,
            fs=16, color=C["muted"], align="center")

# 기술 스택 칩
rect(s, 0.55, 4.88, W - 1.1, 0.6, C["tagBg"], "253E6A", 0.7)
txt(s, "기술 스택 : ", 0.75, 4.92, 2.0, 0.52, fs=12, bold=True, color=C["muted"])
chips = [
    ("Python 3.11",  C["accent2"]),
    ("Streamlit",    C["accent1"]),
    ("Folium",       C["accent3"]),
    ("Pandas",       C["warn"]),
    ("requests",     C["orange"]),
    ("Shapely",      C["danger"]),
    ("Geopy",        C["pink"]),
    ("Haversine",    "7C3AED"),
]
cx = 2.75
for label, color in chips:
    cw = len(label) * 0.1 + 0.9
    rect(s, cx, 4.98, cw, 0.4, C["cardBg"], color, 1.0)
    txt(s, label, cx, 4.98, cw, 0.4, fs=10, bold=True, color=color, align="center")
    cx += cw + 0.18

# 최종 결과물
rect(s, 0.55, 5.62, W - 1.1, 1.52, C["cardBg"], "253E6A", 0.7)
txt(s, "📦  최종 산출물", 0.75, 5.7, 3.5, 0.48,
    fs=13, bold=True, color=C["accent1"])
deliverables = [
    "data_processor.py — 데이터 로딩·전처리·지수 계산 모듈",
    "app.py — Streamlit 웹 애플리케이션  (지도 + 검색 + 테이블)",
]
for i, d in enumerate(deliverables):
    txt(s, f"  ›  {d}", 0.75, 6.25 + i * 0.45, W - 1.3, 0.42,
        fs=11.5, color=C["offWhite"], font="Consolas")


# ════════════════════════════════════════════════════════
# SLIDE 11: 기대 효과 & 의의
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["contentBg"])
header_bar(s, "기대 효과 & 프로젝트 의의", "이 서비스가 왜 필요한가")
step_badge(s, "STEP 08  기대 효과", C["accent3"])
page_num(s, 11)

effects = [
    ("🙋", "전기차 구매\n예정자",
     "내 동네의 충전 인프라 현황을 한눈에 파악\n구매 결정에 실질적 도움",
     C["accent1"]),
    ("🚗", "기존 전기차\n보유자",
     "현재 위치에서 가장 가까운 충전소 즉시 검색\n충전 불안 해소",
     C["accent3"]),
    ("🏛️", "서울시·정책\n담당자",
     "인프라 부족 지역 데이터 기반 파악\n충전소 설치 우선순위 결정에 활용",
     C["warn"]),
    ("📚", "데이터 분석\n학습 관점",
     "공공데이터 수집·전처리·시각화·웹 서비스\n전 과정을 파이썬으로 구현하는 실습",
     C["pink"]),
]
ew = (W - 1.1) / 4
for i, (icon, who, desc, color) in enumerate(effects):
    x = 0.55 + i * (ew + 0.02)
    rect(s, x, 1.6, ew, 4.0, C["midBg"], "253E6A", 0.7)
    rect(s, x, 1.6, ew, 0.09, color)
    oval(s, x + ew/2 - 0.45, 1.82, 0.9, 0.9, color)
    txt(s, icon, x + ew/2 - 0.45, 1.82, 0.9, 0.9,
        fs=24, align="center", valign="middle", color=C["white"])
    txt(s, who, x + 0.12, 2.85, ew - 0.25, 0.7,
        fs=12.5, bold=True, color=C["white"], align="center")
    txt(s, desc, x + 0.12, 3.62, ew - 0.25, 1.75,
        fs=11, color=C["offWhite"], align="center")

# 핵심 문구
rect(s, 0.55, 5.82, W - 1.1, 1.32, C["tagBg"], C["accent3"], 1.0)
txt(s,
    '"공공데이터 + 파이썬으로 시민 생활에 실질적 가치를 제공하는 서비스를 만든다"',
    0.75, 5.88, W - 1.3, 1.2,
    fs=16, bold=True, italic=True, color=C["accent3"], align="center")


# ════════════════════════════════════════════════════════
# SLIDE 12: Q&A
# ════════════════════════════════════════════════════════
s = add_slide()
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(C["darkBg"])

oval(s, -1.5, 4.0, 6.5, 6.5, "001530")
oval(s, 10.5, -1.5, 5.5, 5.5, "001A3A")
rect(s, 0, 0, 0.35, H, C["accent1"])

txt(s, "Q & A", 1.0, 1.2, W - 2.0, 2.2,
    fs=78, bold=True, color=C["white"], align="center")
rect(s, 3.0, 3.95, 7.33, 0.055, "253E6A")

txt(s, "감사합니다", 1.0, 4.18, W - 2.0, 0.88,
    fs=34, color=C["offWhite"], align="center")
txt(s,
    "서울시 전기차 충전 인프라 부족 지역 분석 및 시각화 서비스  |  주제 제안 발표",
    1.0, 5.18, W - 2.0, 0.5,
    fs=13, italic=True, color=C["muted"], align="center")

rect(s, 4.2, 6.05, 4.93, 0.9, C["tagBg"], "253E6A", 0.8)
txt(s, "팀명 / 발표자 이름 입력", 4.2, 6.05, 4.93, 0.9,
    fs=13, italic=True, color=C["muted"], align="center", valign="middle")

page_num(s, 12)


# ─── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT_FILE)
print(f"✅ 주제 제안 PPT 생성 완료: {OUT_FILE}")
print(f"   슬라이드 수: {TOTAL}장")
